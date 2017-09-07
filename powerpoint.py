# -*- coding: utf-8 -*-
from pptx import Presentation
from datetime import datetime
import os
import shutil
import re
from utils import *
from operator import itemgetter
import pythoncom
import logging
from imgur import get_image
import threading
import urllib

def build_report_slide(prs, missionary, report, report_split):
    # Access placeholders for content slides
    content_slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Add biography, mainly from factfile
    insert_bio(content_slide, missionary, report)

    # Report title - pull report round out of date
    enter_report_title(report, content_slide)

    # Actual report!
    report_holder = content_slide.placeholders[14]
    report_holder.text_frame.clear()
    p = report_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report_split.rstrip().rstrip(u"\u2202")

    # Prayer heading
    prayer_h_holder = content_slide.placeholders[15]
    prayer_h_holder.text_frame.clear()
    p = prayer_h_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Prayer Points"

    # TODO - Combine prayers in dictionary together into Prayer points body
    prayer_b_holder = content_slide.placeholders[16]
    prayer_b_holder.text_frame.clear()
    p = prayer_b_holder.text_frame.paragraphs[0]
    run = p.add_run()
    [prayer.capitalize() for prayer in report.prayer_rqs]
    run.text = "\n".join(report.prayer_rqs)

    logging.info("Added report slide for {0}:{1}".format(missionary.id,
                                                         report.round))
    return


def insert_bio(slide, missionary, report):
    logging.info("Inserting bio for {0}".format(missionary.id))
    # Insert state
    state_holder = slide.placeholders[2]
    state_holder.text_frame.clear()
    p = state_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = "State: " + missionary.state
    except AttributeError:
        missionary.state = validate_state(missionary.id[:2],
                                          abbreviation=True,
                                          convert_to_full=True)
        run.text = "State: " + missionary.state

    # Insert India Map based off state name
    india_pic_holder = slide.placeholders[12]
    try:
        india_pic_holder.insert_picture('C:\Users\\br1\Dropbox\NCM\Reports' +
                                        '\!Reporting Workflow\Map Images\\' +
                                        missionary.state + '.png')
    except IOError:
        logging.error("Missing state map for {0}, not added".format(
            missionary.state))

    # Insert Name
    name_holder = slide.placeholders[13]
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = missionary.first_name + " " + missionary.surname
    except AttributeError:
        run.text = missionary.surname

    bio_holder = slide.placeholders[11]
    bio_holder.text_frame.clear()
    p = bio_holder.text_frame.paragraphs[0]
    # Apply numbers for churches and baptisms
    churches = 0
    prayer_nos = 0
    baptisms = 0

    for village in report.villages:
        churches += 1
        if village.attendance:
            try:
                prayer_nos += int(village.attendance)
            except ValueError:
                logging.info("Non-numerical value for attendance")
                try:
                    prayer_nos += int(village.attendance.split("-")[1])
                except ValueError:
                    logging.error("Invalid value for attendance {0}, "
                                  "Missionary ID: {1}".format(
                                      village.attendance,
                                      missionary.id))
        if village.baptisms:
            try:
                baptisms += int(village.baptisms)
            except ValueError:
                logging.error("Invalid value for baptisms {0}, "
                              "Missionary ID: {1}".format(
                                  village.baptisms,
                                  missionary.id))

    bio_line("\n Churches: ", str(churches), p)
    bio_line("\n Coming for Prayer: ", str(prayer_nos), p)
    bio_line("\n Baptisms: ", str(baptisms), p)

    # Download picture to offline, store off and add to report
    logging.info("Adding profile pic for {0}".format(missionary.id))
    profile_pic_holder = slide.placeholders[10]
    try:
        img_filename = missionary.pic.split('/')[-1].rstrip("\'")
        if img_filename:
            urllib.urlretrieve(missionary.pic, img_filename)
            profile_pic_holder.insert_picture(img_filename)
            os.remove(img_filename)
        else:
            raise AttributeError
    except (AttributeError, IOError) as e:
        logging.warning("No headshot for {0}.".format(missionary.id))
        profile_pic_holder.insert_picture(
            "C:\Users\\br1\Dropbox\NCM\Reports\Ben Report Automation" +
            "\headshot.png")
    except Exception as e:
        logging.exception("Profile pic error for {0}".format(missionary.id))

    logging.info("Bio inserted for {0}".format(missionary.id))


def enter_report_title(report, slide):
    title_holder = slide.placeholders[0]
    title_holder.text_frame.clear()
    p = title_holder.text_frame.paragraphs[0]
    run = p.add_run()
    if report.round:
        run.text = str(report.round[1]) + " Report " + str(report.round[0])
    else:
        run.text = "Report"


def create_title_slide(prs, missionary):
    # Access placeholders for Title slide
    title_slide = prs.slides[0]

    # Insert Missionary Name
    name_holder = title_slide.placeholders[0]
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = missionary.first_name + " " + missionary.surname
    except AttributeError:
        run.text = missionary.surname

    # Insert Missionary ID
    miss_id_holder = title_slide.placeholders[11]
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Missionary ID: " + missionary.id

    # Insert current year
    miss_id_holder = title_slide.placeholders[1]
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Reports - " + str(datetime.now().year)

    logging.info("Title slide for {0} complete".format(missionary.id))
    return prs


def create_powerpoint(missionary):
    # Import presentation
    path = "C:\Users\\br1\Dropbox\NCM\Reports\Ben Report Automation\\"
    prs = Presentation(path + "Master Report Template.pptx")

    # Title slide - requires name and ID only.
    create_title_slide(prs, missionary)

    counter = 1
    for report_no, report in sorted(sorted(missionary.reports.iteritems(),
                                           key=itemgetter(0),
                                           reverse=True),
                                    key=itemgetter(1),
                                    reverse=True):
        logging.info("Creating report slide for {0} {1}".format(
            missionary.id, report.round))
        try:
            for report_split in report.report:
                build_report_slide(prs, missionary, report, report_split)
                counter += 1
        except AttributeError:
            logging.error("No report for {0}, reports: {1}".format(
                missionary.id, missionary.reports))

    # Save the powerpoint in a folder with Missionary ID
    path = "C:\Users\\br1\Code\\500k\\reports\\{0}_{1}.pptx".format(
        missionary.id,
        missionary.first_name)
    prs.save(path)
    logging.info("Reports for {0} have been saved to {1}.".format(
                    missionary.id, path))
    return path


def create_powerpoint_pdf(q):
    while True:
        try:
            logging.info("Thread engaged.")
            (missionary, miss_id, date) = q.get()
            create = False
            if date:
                month = int(date.split('/')[0])
                year = int(date.split('/')[1])
                for k,v in missionary.reports.iteritems():
                    if v.get_month() == month and v.get_year() == year:
                        logging.info("Found report in month {0} for {1}".format(
                            month,
                            miss_id))
                        create = True
                        break
            else:
                logging.info("No date provided, make for {0}".format(miss_id))
                create = True

            if create and missionary.reports:
                logging.info("Building pptx file for {0}".format(miss_id))
                path = create_powerpoint(missionary)

                pythoncom.CoInitialize()
                # Export to pdf - this is the slowest part so thread.
                if path:
                    pdf_path = path.split(".")[0] + ".pdf"
                    try:
                        PPTtoPDF(path, pdf_path)
                        logging.info("Converted to PDF: {0}".format(pdf_path))
                        try:
                            os.remove(path)
                        except OSError:
                            logging.warning("Could not find {0} to delete".format(
                                path))
                    except Exception as e:
                        logging.exception(
                            "Build PDF failed for missionary with ID: {0}".format(
                                miss_id))
                else:
                    logging.error(
                        "Missing pptx for missionary with ID: {0}".format(
                            miss_id))
        except Exception as e:
            logging.exception("{0} has died!".format(
                          threading.current_thread().name))
        finally:
            logging.info("Releasing thread for next job.")
            q.task_done()

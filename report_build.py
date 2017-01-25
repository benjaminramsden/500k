from pptx import Presentation
from docx import Document
from imgurpython import ImgurClient
from datetime import datetime
import sys
from utils import copy_unzip_docx, find_pic_in_docx, bio_line

def main(argv=None):
    # First attempt to get all the information we need out of the word doc
    print "What's the path to the report docx we are using?"
    docx_path = raw_input()
    if not docx_path:
        docx_path = "C:\Users\\br1\Dropbox\NCM\Reports, bills and Proposals\Ben Report Automation\AP 16 12  M Raj Kumar Rajan KL1234.docx"
    doc = Document(docx_path)

    # Get the round and year out of the report filename - TODO
    # The format of this is not clear yet...

    # Time for info gathering, info we need:
    # - Name
    # - State
    # - Numbers (if any)
    # - Report itself
    # - Prayer Points

    # Standard bio info
    for para in doc.paragraphs:
        if para.text.startswith("Name") and not para.text.startswith("Name of"):
            name = para.text.split(":")[-1][1:]
        elif para.text.startswith("Date of Birth"):
            dob = para.text.split(":")[-1][1:].strip(" ")
            d = datetime.strptime(dob, '%d/%m/%Y')
            age = str(datetime.now().year - d.year)
        elif para.text.startswith("Wife"):
            wife = para.text.split(":")[-1][1:]
            if wife.isspace():
                wife = "None"
        elif para.text.startswith("Children"):
            children = para.text.split(":")[-1][1:]
            if children.isspace():
                children = "None"
        elif para.text.startswith("Languages"):
            languages = para.text.split(":")[-1][1:]
            if languages.isspace():
                languages = "None"

    # Info on number of churches and baptisms
    churches = 0
    baptisms = 0
    prayer_nos = 0

    for cell in doc.tables[0].column_cells(3)[1:]:
        if cell.text:
            baptisms += int(cell.text)
            churches += 1

    for cell in doc.tables[0].column_cells(2)[1:]:
        if cell.text:
            prayer_nos += int(cell.text)

    # Finally retrieve the report and prayer points
    report = ""
    for idx, para in enumerate(doc.paragraphs):
        if para.text in ["Prayer Requests", "Prayer Points"]:
            for para in doc.paragraphs[:idx]:
                if len(para.text) > 150:
                    report += para.text + "\n"

    # Look at parent directory strip the state from the first two characters
    # in the ID. Create a dictionary of two-letter abbreviations to state names
    parent_dir = docx_path.split("\\")[-2]
    miss_id = parent_dir[-6]
    state = parent_dir[-6:-4]

    # Use a dictionary to convert the two-letter state acronym to full name - TODO

    # Prayer points are hard with the bullets
    prayer = ""
    for idx, para in enumerate(doc.paragraphs):
        if para.text in ["Prayer Requests", "Prayer Points"]:
            for line in doc.paragraphs[idx+1:]:
                if line.text:
                    prayer += line.text + "\n"

    # Remove trailing return so we don't get extra bullet point
    prayer.rstrip("\n")

    # Don't forget their profile picture! Get this by unzipping the file
    unzip_path = copy_unzip_docx(docx_path)
    img_path = find_pic_in_docx(unzip_path)

    print "Enter path for master presentation template (will use Ben's default if blank)"
    path = raw_input()
    if not path:
        path = "C:\Users\\br1\Dropbox\NCM\Reports, bills and Proposals\Ben Report Automation\\"

    # Import presentation
    prs = Presentation(path + "Master Report Template.pptx")

    # Add extra report slides based on available number
    print "How many reports are there to put together?"
    report_no = raw_input()

    if int(report_no) > 1:
        # Create new slides for each other report to add, make sure to insert same
        # placeholders
        new_slides = report_no -1

        # Not tested!!
        # prs.slides.add_slide(slide_layouts[0])

    # Access placeholders for Title slide

    title_slide = prs.slides[0]

    for shape in title_slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

    # Insert Missionary Name
    name_holder = title_slide.placeholders[0]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = name

    # Insert Missionary ID - TODO

    # Access placeholders for content slides
    content_slide = prs.slides[1]

    # Insert the name and the state of the Missionary
    state_holder = content_slide.placeholders[2]
    assert state_holder.has_text_frame
    state_holder.text_frame.clear()
    p = state_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "State: " + state

    name_holder = content_slide.placeholders[13]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = name

    bio_holder = content_slide.placeholders[11]
    assert bio_holder.has_text_frame
    bio_holder.text_frame.clear()
    p = bio_holder.text_frame.paragraphs[0]
    bio_line("Age: ", age, p)
    bio_line("\n Spouse: ", wife, p)
    bio_line("\n Children: ", children, p)
    bio_line("\n Languages: ", languages, p)

    profile_pic_holder = content_slide.placeholders[10]

    # Insert profile picture, need to access Imgur database!
    profile_pic_holder.insert_picture(img_path)

    # Insert India Map based off state name
    #india_pic_holder = content_slide.placeholders[12]
    #india_pic_holder.insert_picture('C:\Users\\br1\Dropbox\NCM\Reports, bills and Proposals\Ben Report Automation\Map images\\' + state.lower() + '.png')

    # Report title
    title_holder = content_slide.placeholders[0]
    assert title_holder.has_text_frame
    title_holder.text_frame.clear()
    p = title_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Report - <Year> Report <round>"

    # Actual report!
    report_holder = content_slide.placeholders[14]
    assert report_holder.has_text_frame
    report_holder.text_frame.clear()
    p = report_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report

    # Prayer heading
    prayer_h_holder = content_slide.placeholders[15]
    assert prayer_h_holder.has_text_frame
    prayer_h_holder.text_frame.clear()
    p = prayer_h_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Prayer Points"

    # Prayer points body
    prayer_b_holder = content_slide.placeholders[16]
    assert prayer_b_holder.has_text_frame
    prayer_b_holder.text_frame.clear()
    p = prayer_b_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = prayer

    # Save the powerpoint
    prs.save('test.pptx')

    # Tidy up unzipped word doc and .zip file - TODO
    try:
        pass
    finally:
        pass

if __name__ == '__main__':
    status = main()
    sys.exit(status)

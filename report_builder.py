# -*- coding: utf-8 -*-
from pptx import Presentation
from docx import Document
from datetime import datetime
import sys
import os
import shutil
import re
from utils import *
from sheets_api import *
from report import Report
from village import Village
from missionary import Missionary, Child, Spouse
from operator import itemgetter
from Queue import Queue
import threading
import pythoncom
import logging
from imgur import get_image

# This script conducts the following:
# - Gets the information on a missionary based on Miss ID (gets all)
# - Populates the title slide with the missionaries details
# - Creates a content slide per report
# - Exports the slideshow to pdf
# - Uploads to Ben's Google Drive account via API
# - Pastes the URL of the report in Google Drive to the web report sheets
#
# With this info Ben can then send out multiple reports using:
# https://support.yet-another-mail-merge.com/hc/en-us/articles/210735349


def main(argv=None):
    logging.basicConfig(filename='diags.log',
                        level=logging.WARNING,
                        format='%(asctime)s %(name)-12s %(levelname)-8s %(message)s',
                        datefmt='%m-%d %H:%M')

    # Gather all information from the spreadsheet. Returned as list of lists
    # where each list is a row of cells.
    report_data = get_all_missionary_reports(test=True)

    # Add in the factfile information
    factfile_data = get_all_factfile_data()

    # Now build out the data into usable dictionaries
    all_missionaries = construct_data(report_data, factfile_data)

    # Time to create the presentations, loop around for every single missionary
    # TODO - In future make sure only missionaries with new reports get
    # generated
    logging.info("Creating powerpoints for {0} missionaries".format(
        len(all_missionaries)))

    q = Queue(maxsize=0)
    num_threads = 10

    for i in range(num_threads):
        worker = threading.Thread(target=create_powerpoint_pdf, args=(q,))
        worker.setDaemon(True)
        worker.start()

    for miss_id, missionary in all_missionaries.iteritems():
        q.put((missionary, miss_id))

    q.join()

    # TODO - Now upload all these reports to Google Drive via API, saving the
    # URL/ID of the report back into Google Sheets

    return 0


def construct_data(report_data, factfile_data):
    """
    Take from the two different spreadsheets to create a total view of all the
    missionary data, once complete we have all the info required to start
    creating the reports.
    """
    all_missionaries = {}
    construct_factfile_data(all_missionaries, factfile_data)
    construct_report_data(all_missionaries, report_data)
    return all_missionaries


def construct_report_data(all_missionaries, report_data):
    # For all the missionaries, arrange data in this structure:
    # All
    #  -> Missionary 1 (based on ID)
    #      -> Report 1
    #          -> Date, Subject, Raw, Missionary, Missionary ID etc...
    #          -> Village 1
    #              -> Village
    #              -> People
    #              -> Baptisms
    #          -> Village 2
    #           ...
    #          -> All Prayer Points
    #           ...
    #      -> Report 2
    #       ...
    #  -> Missionary 2
    #   ...
    logging.info("Constructing report data")

    # As we may change the order of the columns from time to time and need to
    # make this sustainable for any changes, create a dictionary of column
    # numbers against the column header text. This should be a single linear
    # search to get all the headings.
    columns = dict()
    for idx, column in enumerate(report_data[0]):
        columns[column] = idx

    for row in report_data[1:]:
        if len(row) > columns[u'\u2022Main Story / Report: ']:
            try:
                report = Report(row[columns['Date (Pretty)']],
                                row[columns[u'\u2022Missionary Name: ']],
                                row[columns[u'\u2022MissionaryID: ']],
                                row[columns[u'\u2022Main Story / Report: ']])
            except NotImplementedError:
                continue
            villages = []
            prayer_rqs = []
            for i in range(1, 6):
                if row[columns[u'\u2022V' + str(i) + ': ']]:
                    villages.append(
                        Village(row[columns[u'\u2022V' + str(i) + ': ']],
                                row[columns[u'\u2022V' + str(i) + 'N: ']],
                                row[columns[u'\u2022V' + str(i) + 'B: ']]))
            for i in range(1, 9):
                if (len(row) > columns["P-R-" + str(i) + ": "] and
                        row[columns["P-R-" + str(i) + ": "]]):
                    prayer_rqs.append(row[columns["P-R-" + str(i) + ": "]])
            report.villages = villages
            report.prayer_rqs = prayer_rqs
            report.round = report.get_report_round()
            missionary_id = report.id
            if missionary_id in all_missionaries.keys():
                # Missionary already exists, add report to dictionary
                missionary = all_missionaries[missionary_id]
            else:
                # New missionary, create new missionary and add report.
                logging.warning("Missionary not found, does {0} exist?".format(
                    missionary_id))
                names = report.name.split(" ")
                if len(names) > 1:
                    try:
                        missionary = Missionary(missionary_id,
                                                names[-1],
                                                names[-2])
                    except NotImplementedError:
                        continue
                else:
                    try:
                        missionary = Missionary(missionary_id,
                                                names[-1],
                                                None)
                    except NotImplementedError:
                        continue
                all_missionaries[missionary_id] = missionary
            missionary.reports[report.round] = report

    logging.info("Report data has been constructed")


def construct_factfile_data(all_missionaries, factfile_data):
    """
    With missionary reports constructed, now add the factfile data to the side
    """
    logging.info("Constructing factfile data")

    # As we may change the order of the columns from time to time and need to
    # make this sustainable for any changes, create a dictionary of column
    # numbers against the column header text. This should be a single linear
    # search to get all the headings.
    columns = dict()
    for idx, column in enumerate(factfile_data[0]):
        columns[column] = idx

    for row in factfile_data[1:]:
        if len(row) > columns[u'Profile Picture']:
            # Basics mandatory for a factfile
            try:
                missionary = Missionary(row[columns[u'ID (new)']],
                                        row[columns[u'MissionarySecondName']],
                                        row[columns[u'MissionaryFirstName']])
            except NotImplementedError:
                continue
            try:
                missionary.state = validate_state(
                    row[columns[u'MissionField State']])
            except ValueError:
                continue
            missionary.pic = row[columns[u'Profile Picture']]
            # Add family and biography
            if len(row) > columns[u'Number of Dependents']:
                if row[columns[u'Wife / Husband\'s First Name']]:
                    missionary.spouse = Spouse(
                        row[columns[u'Wife / Husband\'s First Name']],
                        row[columns[u'Wife / Husband\'s Second Name']],
                    )
            for i in range(1, 6):
                if row[columns[u'Child ' + str(i) + ' First Name']]:
                    missionary.children[u'Child ' + str(i)] = Child(
                        row[columns[u'Child ' + str(i) + ' First Name']],
                        row[columns[u'Child ' + str(i) + ' DOB']])

            # Mission field data
            villages = []
            prayer_rqs = []
            for i in range(1, 6):
                if (len(row) > columns[u'V' + str(i) + ' B'] and
                        row[columns[u'V' + str(i)]]):
                    villages.append(
                        Village(row[columns[u'V' + str(i)]],
                                row[columns[u'V' + str(i) + ' N']],
                                row[columns[u'V' + str(i) + ' B']]))
            missionary.villages = villages

    logging.info("Factfile data has been constructed")


def create_powerpoint(missionary):
    # Import presentation
    path = "C:\Users\\br1\Dropbox\NCM\Reports\Ben Report Automation\\"
    prs = Presentation(path + "Master Report Template.pptx")

    # Title slide - requires name and ID only.
    create_title_slide(prs, missionary)

    counter = 1
    for report_no, report in sorted(sorted(missionary.reports.iteritems(),
                                           key=itemgetter(0),
                                           reverse=True),
                                    key=itemgetter(1),
                                    reverse=True):
        for report_split in report.report:
            build_report_slide(prs, missionary, report, report_split)
            counter += 1

    # TODO - Save the powerpoint in a folder with Missionary ID
    path = "C:\Users\\br1\Code\\500k\\reports\\{0}_{1}.pptx".format(
        missionary.id,
        missionary.surname)
    prs.save(path)
    logging.info("Reports for {0} have been saved to {1}.".format(
                    missionary.id, path))
    return path


def create_powerpoint_pdf(q):
    while True:
        try:
            (missionary, miss_id) = q.get()
            path = create_powerpoint(missionary)

            pythoncom.CoInitialize()
            # Export to pdf - this is the slowest part so thread.
            if path:
                try:
                    PPTtoPDF(path, path.split(".")[0] + ".pdf")
                    try:
                        os.remove(path)
                    except OSError:
                        logging.warning("Could not find {0} to delete".format(
                            path))
                except:
                    logging.error(
                        "Build PDF failed for missionary with ID: {0}".format(
                            miss_id))
            else:
                logging.error(
                    "Missing pptx for missionary with ID: {0}".format(
                        miss_id))
        except:
            logging.error("{0} has died!".format(
                          threading.current_thread().name))
        finally:
            q.task_done()


def create_title_slide(prs, missionary):
    # Access placeholders for Title slide
    title_slide = prs.slides[0]

    # Insert Missionary Name
    name_holder = title_slide.placeholders[0]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = missionary.first_name + " " + missionary.surname
    except AttributeError:
        run.text = missionary.surname

    # Insert Missionary ID
    miss_id_holder = title_slide.placeholders[11]
    assert miss_id_holder.has_text_frame
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Missionary ID: " + missionary.id

    # Insert current year
    miss_id_holder = title_slide.placeholders[1]
    assert miss_id_holder.has_text_frame
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Reports - " + str(datetime.now().year)

    return prs


def insert_bio(slide, missionary, report):
    # Insert state
    state_holder = slide.placeholders[2]
    assert state_holder.has_text_frame
    state_holder.text_frame.clear()
    p = state_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = "State: " + missionary.state
    except AttributeError:
        missionary.state = validate_state(missionary.id[:2],
                                          abbreviation=True,
                                          convert_to_full=True)
        run.text = "State: " + missionary.state

    # Insert India Map based off state name
    india_pic_holder = slide.placeholders[12]
    try:
        india_pic_holder.insert_picture('C:\Users\\br1\Dropbox\NCM\Reports' +
                                        '\!Reporting Workflow\Map Images\\' +
                                        missionary.state + '.png')
    except IOError:
        logging.error("ERROR: Missing state map for {0}, not added".format(
            missionary.state))

    # Insert Name
    name_holder = slide.placeholders[13]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = missionary.first_name + " " + missionary.surname
    except AttributeError:
        run.text = missionary.surname

    bio_holder = slide.placeholders[11]
    assert bio_holder.has_text_frame
    bio_holder.text_frame.clear()
    p = bio_holder.text_frame.paragraphs[0]
    # Apply numbers for churches and baptisms
    churches = 0
    prayer_nos = 0
    baptisms = 0

    for village in report.villages:
        churches += 1
        if village.attendance:
            try:
                prayer_nos += int(village.attendance)
            except ValueError:
                logging.info("Non-numerical value for attendance")
                try:
                    prayer_nos += int(village.attendance.split("-")[1])
                except ValueError:
                    logging.error("Invalid value for attendance {0}, "
                                  "Missionary ID: {1}".format(
                                      village.attendance,
                                      missionary.id))
        if village.baptisms:
            try:
                baptisms += int(village.baptisms)
            except ValueError:
                logging.error("Invalid value for baptisms {0}, "
                              "Missionary ID: {1}".format(
                                  village.baptisms,
                                  missionary.id))

    bio_line("\n Churches: ", str(churches), p)
    bio_line("\n Coming for Prayer: ", str(prayer_nos), p)
    bio_line("\n Baptisms: ", str(baptisms), p)

    # Download Imgur picture, store off and add to report
    profile_pic_holder = slide.placeholders[10]
    try:
        if not missionary.pic:
            missionary.pic = get_image(missionary.id)
        profile_pic_holder.insert_picture(missionary.pic)
    except AttributeError:
        logging.error("No headshot for {0}.".format(missionary.id))
        profile_pic_holder.insert_picture(
            "C:\Users\\br1\Dropbox\NCM\Reports\Ben Report Automation" +
            "\headshot.png")

    # get_bio_from_factfile(slide,report["Missionary ID"])
    return


def get_bio_from_factfile(slide, miss_id):
    # Pull down info for missionary based off missionary ID from factfile sheet
    ff_data = get_all_factfile_data()


def enter_report_title(report, slide):
    title_holder = slide.placeholders[0]
    assert title_holder.has_text_frame
    title_holder.text_frame.clear()
    p = title_holder.text_frame.paragraphs[0]
    run = p.add_run()
    if report.round:
        run.text = str(report.round[1]) + " Report " + str(report.round[0])
    else:
        run.text = "Report"


def build_report_slide(prs, missionary, report, report_split):
    # Access placeholders for content slides
    content_slide = prs.slides.add_slide(prs.slide_layouts[0])

    # Add biography, mainly from factfile
    success = insert_bio(content_slide, missionary, report)

    # Report title - pull report round out of date
    enter_report_title(report, content_slide)

    # Actual report!
    report_holder = content_slide.placeholders[14]
    assert report_holder.has_text_frame
    report_holder.text_frame.clear()
    p = report_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report_split.rstrip()

    # Prayer heading
    prayer_h_holder = content_slide.placeholders[15]
    assert prayer_h_holder.has_text_frame
    prayer_h_holder.text_frame.clear()
    p = prayer_h_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Prayer Points"

    # TODO - Combine prayers in dictionary together into Prayer points body
    prayer_b_holder = content_slide.placeholders[16]
    assert prayer_b_holder.has_text_frame
    prayer_b_holder.text_frame.clear()
    p = prayer_b_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "\n".join(report.prayer_rqs)

    return success

if __name__ == '__main__':
    status = main()
    sys.exit(status)

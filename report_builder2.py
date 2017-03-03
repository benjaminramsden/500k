from pptx import Presentation
from docx import Document
from imgurpython import ImgurClient
from datetime import datetime
import sys, os, shutil, re
from utils import *
from sheets_api import *

# This script conducts the following:
# - Gets the information on a missionary based on Miss ID (gets all)
# - Populates the title slide with the missionaries details
# - Creates a content slide per report
# - Exports the slideshow to pdf
# - Uploads to Ben's Google Drive account via API
# - Pastes the URL of the report in Google Drive to the web report sheets
#
# With this info Ben can then send out multiple reports using:
# https://support.yet-another-mail-merge.com/hc/en-us/articles/210735349

def main(argv=None):
    # Gather all information from the spreadsheet. Returned as list of lists
    # where each list is a row of cells.
    values = get_all_missionary_reports()

    all_dict = construct_data(values)

    # Time to create the presentations, loop around for every single missionary
    # TODO - In future make sure only missionaries with new reports get
    # generated
    for miss_id,miss_dict in all_dict.iteritems():
        pptx = create_powerpoint(miss_dict)

        # Export to pdf
        if pptx:
            PPTtoPDF(pptx, pptx.split(".")[0] + ".pdf")
        else:
            print "Build failed for missionary with ID:" + miss_id

    return 0

def construct_data(values):
    # For all the missionaries, arrange data in this structure:
    # All
    #  -> Missionary 1 (based on ID)
    #      -> Report 1
    #          -> Date, Subject, Raw, Missionary, Missionary ID etc...
    #          -> Village 1
    #              -> Village
    #              -> People
    #              -> Baptisms
    #          -> Village 2
    #           ...
    #          -> All Prayer Points
    #           ...
    #      -> Report 2
    #       ...
    #  -> Missionary 2
    #   ...
    print "Constructing data"

    all_dict = dict()
    for row in values:
        report = {"Date":         row[0],
                  "Subject":      row[1],
                  "Raw":          row[3],
                  "Submitter":    row[4],
                  "Email":        row[5],
                  "Missionary":   row[6],
                  "Missionary ID":row[7],
                  "Report":       row[40],
                 }
        for i,village in enumerate(row[8:3:26]):
            if village != "":
                vill_dict = {
                    "Village": row[i+6],
                    "People":  row[i+7],
                    "Baptisms":row[i+8],
                }
                report['Village '+str(i+1)] = vill_dict
        report['Prayer'] = '\n'.join(row[41:48]).strip()
        report_round = get_report_round(report["Date"])
        if report["Missionary ID"] in all_dict.keys():
            # Missionary already exists, add report to missionary dictionary
            miss_dict = all_dict[report["Missionary ID"]]
            miss_dict[report_round] = report
        else:
            # New missionary, create new dictionary and add report to it.
            all_dict[report["Missionary ID"]] = {report_round: report}

    print "Data has been constructed"
    return all_dict

def create_powerpoint(miss_dict):
    # Import presentation
    print "Opening powerpoint template"
    path = "C:\Users\\br1\Dropbox\NCM\Reports\Ben Report Automation\\"
    prs = Presentation(path + "Master Report Template.pptx")

    # Grab the first report in the dictionary for this info, it doesn't matter
    # which
    create_title_slide(prs,miss_dict.itervalues().next())

    counter = 1
    for report in miss_dict.itervalues():
        print "Creating report slide " + str(counter)
        build_report_slide(prs,report)
        counter+=1

    # TODO - Save the powerpoint in a folder with Missionary ID - discuss with Alex
    path = "C:\Users\\br1\Code\\500k\\{0}_{1}.pptx".format(
        report["Missionary ID"],
        report["Missionary"])
    prs.save(path)
    return path

def create_title_slide(prs,report):
    # Access placeholders for Title slide
    title_slide = prs.slides[0]

    # Insert Missionary Name
    name_holder = title_slide.placeholders[0]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report["Missionary"]

    # Insert Missionary ID
    miss_id_holder = title_slide.placeholders[11]
    assert miss_id_holder.has_text_frame
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Missionary ID: " + report["Missionary ID"]

    # Insert current year
    miss_id_holder = title_slide.placeholders[1]
    assert miss_id_holder.has_text_frame
    miss_id_holder.text_frame.clear()
    p = miss_id_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Reports - " + str(datetime.now().year)

    return prs

def insert_bio(slide,report):
    # TODO - Get totals for churches, baptisms and prayer
    #for k, value in report:
    #    # Only villages are nested dictionaries, so test on value is dictionary
    #    # BLEURGH
    #    if isinstance(value, dict):
    #        pass

    # Define the state acronyms here, this could be moved out at a later date
    state_dict = {"AN": "Andaman Nicobar",
                  "AP": "Andhra Pradesh",
                  "AS": "Assam",
                  "CH": "Chhattisgarh",
                  "HR": "Haryana",
                  "GJ": "Gujarat",
                  "HP": "Himachal Pradesh",
                  "JK": "Jammu & Kashmir",
                  "KA": "Karnataka",
                  "KL": "Kerala",
                  "MP": "Madhya Pradesh",
                  "MH": "Maharashtra",
                  "OR": "Odisha (Orissa)",
                  "PB": "Punjab",
                  "RJ": "Rajasthan",
                  "TN": "Tamil Nadu",
                  "TA": "Telangana",
                  "TR": "Tripura",
                  "UP": "Uttar Pradesh",
                  "UK": "Uttarakhand"}

    state_ab = report["Missionary ID"][:2]
    # Insert state
    state_holder = slide.placeholders[2]
    assert state_holder.has_text_frame
    state_holder.text_frame.clear()
    p = state_holder.text_frame.paragraphs[0]
    run = p.add_run()
    try:
        run.text = "State: " + state_dict[state_ab]
    except KeyError:
        print("Report dated {0} for missionary {1} has incorrect Missionary "
            "ID, chase Shijo".format(report["Date"],report["Missionary"]))
        return 1

    # Insert India Map based off state name
    india_pic_holder = slide.placeholders[12]
    india_pic_holder.insert_picture('C:\Users\\br1\Dropbox\NCM\Reports' +
        '\!Reporting Workflow\Map Images\\' +
        state_dict[state_ab] + '.png')

    # Insert Name
    name_holder = slide.placeholders[13]
    assert name_holder.has_text_frame
    name_holder.text_frame.clear()
    p = name_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report["Missionary"]

    bio_holder = slide.placeholders[11]
    assert bio_holder.has_text_frame
    bio_holder.text_frame.clear()
    p = bio_holder.text_frame.paragraphs[0]
    # TODO - get these numbers out above
    #bio_line("\n Churches: ", str(churches), p)
    #bio_line("\n Coming for Prayer: ", str(prayer_nos), p)
    #bio_line("\n Baptisms: ", str(baptisms), p)

    profile_pic_holder = slide.placeholders[10]

    # get_bio_from_factfile(slide,report["Missionary ID"])
    return

def get_bio_from_factfile(slide,miss_id):
    # Pull down info for missionary based off missionary ID from factfile sheet
    ff_data = get_all_factfile_data()

    # Get the Imgur image ID for the profile picture

    # Pull down Imgur profile picture and insert into presentation
    profile_pic_holder.insert_picture(img_path)

def enter_report_title(report, slide):
    date = datetime.strptime(report["Date"], '%Y-%m-%d %H:%M:%S')
    year = str(date.year)
    report_round = get_report_round(report["Date"])
    title_holder = slide.placeholders[0]
    assert title_holder.has_text_frame
    title_holder.text_frame.clear()
    p = title_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = year + " Report " + report_round.split("/")[0]

def get_report_round(timestamp):
    date = datetime.strptime(timestamp, '%Y-%m-%d %H:%M:%S')
    if date.month in range(1,5):
        report_round = "1/"+str(date.year)
    elif date.month in range(5,9):
        report_round = "2/"+str(date.year)
    elif date.month in range(9,13):
        report_round = "3/"+str(date.year)
    else:
        report_round = "ERROR"
    return report_round

def build_report_slide(prs,report):
    # Access placeholders for content slides
    content_slide = prs.slides.add_slide(prs.slide_layouts[0])

    for shape in content_slide.placeholders:
        print('%d %s' % (shape.placeholder_format.idx, shape.name))

    # Add biography, mainly from factfile
    success = insert_bio(content_slide, report)

    # Report title - pull report round out of date
    enter_report_title(report,content_slide)

    # Actual report! Tidy up regex hacks
    body = report["Report"].replace(">>> ","\n")
    body = re.sub('[ \t\f\v+]', ' ', body)
    report_holder = content_slide.placeholders[14]
    assert report_holder.has_text_frame
    report_holder.text_frame.clear()
    p = report_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = body.rstrip()

    # Prayer heading
    prayer_h_holder = content_slide.placeholders[15]
    assert prayer_h_holder.has_text_frame
    prayer_h_holder.text_frame.clear()
    p = prayer_h_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Prayer Points"

    # TODO - Combine prayers in dictionary together into Prayer points body
    prayer_b_holder = content_slide.placeholders[16]
    assert prayer_b_holder.has_text_frame
    prayer_b_holder.text_frame.clear()
    p = prayer_b_holder.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = report["Prayer"]

    return success

if __name__ == '__main__':
    status = main()
    sys.exit(status)
